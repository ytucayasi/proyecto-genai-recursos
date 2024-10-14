import os
import openai
from pptx import Presentation
import random
import re
import subprocess
from fastapi import HTTPException
from typing import Dict
from dotenv import load_dotenv

# Cargar las variables de entorno
load_dotenv()

# Configurar la API key de OpenAI
openai.api_key = os.getenv("OPENAI_API_KEY")

# Prompt para generar la presentación
Prompt = """Write a presentation/powerpoint about the user's topic. You only answer with the presentation. Follow the structure of the example.
Notice
-You do all the presentation text for the user.
-You write the texts no longer than 250 characters!
-You make very short titles!
-You make the presentation easy to understand.
-The presentation has a table of contents.
-The presentation has a summary.
-At least 8 slides.

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
...

#Slide: 2
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 3
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 5
#Headers: summary
#Content: CONTENT OF THE SUMMARY

#Slide: END"""

def create_ppt_text(input_text: str):
    try:
        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": Prompt},
                {"role": "user", "content": "The user wants a presentation about " + input_text}
            ],
            temperature=0.5,
        )
        return response.choices[0].message.content
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

def create_ppt(text_file: str, design_number: int, ppt_name: str, requested_slide_count: int) -> Dict[str, str]:
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1

    # Eliminar la extensión .txt del nombre del archivo si está presente
    ppt_name = os.path.splitext(ppt_name)[0]

    try:
        with open(text_file, 'r', encoding='utf-8') as f:
            for line in f:
                if slide_count >= requested_slide_count:
                    break

                if line.startswith('#Title:'):
                    header = line.replace('#Title:', '').strip()
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    slide.shapes.title.text = header
                    slide_count += 1
                elif line.startswith('#Slide:'):
                    if slide_count > 0:
                        slide = prs.slides.add_slide(prs.slide_layouts[last_slide_layout_index])
                        slide.shapes.title.text = header
                        body_shape = slide.shapes.placeholders[1]
                        tf = body_shape.text_frame
                        tf.text = content
                        slide_count += 1
                    content = ""
                    last_slide_layout_index = random.choice([1, 7, 8, 9, 10])
                elif line.startswith('#Header:'):
                    header = line.replace('#Header:', '').strip()
                elif line.startswith('#Content:'):
                    content = line.replace('#Content:', '').strip()
                    next_line = f.readline().strip()
                    while next_line and not next_line.startswith('#'):
                        content += '\n' + next_line
                        next_line = f.readline().strip()

        ppt_dir = 'GeneratedPresentations'
        pdf_dir = 'GeneratedPdf'

        os.makedirs(ppt_dir, exist_ok=True)
        os.makedirs(pdf_dir, exist_ok=True)

        ppt_path = os.path.join(ppt_dir, f'{ppt_name}.pptx')
        prs.save(ppt_path)

        if not os.path.exists(ppt_path):
            raise HTTPException(status_code=500, detail="PPTX file not found")

        # Convertir PPTX a PDF
        pdf_path = os.path.join(pdf_dir, f'{ppt_name}.pdf')
        
        # Buscar LibreOffice en ubicaciones comunes
        libreoffice_paths = [
            'libreoffice',  # Si está en el PATH
            '/usr/bin/libreoffice',
            '/usr/local/bin/libreoffice',
            'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
            'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe',
        ]

        conversion_successful = False
        for libreoffice_path in libreoffice_paths:
            try:
                result = subprocess.run(
                    [libreoffice_path, '--headless', '--convert-to', 'pdf', '--outdir', pdf_dir, ppt_path],
                    capture_output=True, 
                    text=True, 
                    check=True
                )
                conversion_successful = True
                break
            except subprocess.CalledProcessError as e:
                print(f"Error with {libreoffice_path}: {e}")
            except FileNotFoundError:
                print(f"LibreOffice not found at {libreoffice_path}")

        if not conversion_successful:
            raise HTTPException(status_code=500, detail="Failed to convert PPTX to PDF. Please ensure LibreOffice is installed and accessible.")

        if not os.path.exists(pdf_path):
            raise HTTPException(status_code=500, detail="PDF file not found after conversion")

        return {
            "pptx": f"/GeneratedPresentations/{ppt_name}.pptx",
            "pdf": f"/GeneratedPdf/{ppt_name}.pdf"
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error creating presentation: {str(e)}")

def generate_filename(input_string: str) -> str:
    filename_prompt = f"Generate a short, descriptive filename based on the following input: \"{input_string}\". Answer just with the short filename, no other explanation."
    try:
        filename_response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": filename_prompt},
            ],
            temperature=0.5,
            max_tokens=30,
        )
        return filename_response.choices[0].message.content.strip().replace(" ", "_").replace(".pptx", "")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))